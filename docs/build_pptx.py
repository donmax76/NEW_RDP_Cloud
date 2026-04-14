#!/usr/bin/env python3
"""Generate Prometey AZ presentation with dark theme, screenshots, and logo."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

IMG = os.path.join(os.path.dirname(__file__), "img")
OUT = os.path.join(os.path.dirname(__file__), "Prometey_AZ.pptx")

BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x00, 0xA8, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x88, 0x88, 0x88)
DARK2 = RGBColor(0x2E, 0x40, 0x57)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_logo(slide, x=Inches(0.3), y=Inches(0.2), h=Inches(0.5)):
    logo = os.path.join(IMG, "logo.jpg")
    if os.path.exists(logo):
        img = Image.open(logo)
        aspect = img.size[0] / img.size[1]
        w = int(h * aspect)
        slide.shapes.add_picture(logo, x, y, w, h)

def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(SW - Inches(0.8), SH - Inches(0.5), Inches(0.6), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT

def add_title(slide, text, x=Inches(1.1), y=Inches(0.15), w=Inches(11), size=28):
    txBox = slide.shapes.add_textbox(x, y, w, Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = ACCENT

def add_text(slide, text, x=Inches(0.5), y=Inches(1.0), w=Inches(5.5), h=Inches(5.5), size=14):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)
    return txBox

def add_image(slide, fname, x, y, max_w, max_h=None):
    path = os.path.join(IMG, fname)
    if not os.path.exists(path):
        return
    img = Image.open(path)
    iw, ih = img.size
    scale = min(max_w / iw, 1.0)
    if max_h:
        scale = min(scale, max_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    slide.shapes.add_picture(path, Emu(x), Emu(y), Emu(w), Emu(h))

def img_slide(num, title, fname, desc):
    """Standard slide: title + description top, screenshot centered below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    add_logo(slide)
    add_title(slide, title)
    add_text(slide, desc, x=Inches(0.5), y=Inches(0.9), w=Inches(12), h=Inches(0.8), size=13)

    path = os.path.join(IMG, fname)
    if os.path.exists(path):
        img = Image.open(path)
        iw_px, ih_px = img.size
        # Convert pixels to EMU (1 inch = 96 pixels = 914400 EMU)
        emu_per_px = 914400 / 96
        iw_emu = int(iw_px * emu_per_px)
        ih_emu = int(ih_px * emu_per_px)
        max_w = Inches(11)
        max_h = Inches(5.0)
        scale = min(max_w / iw_emu, max_h / ih_emu, 1.0)
        pw = int(iw_emu * scale)
        ph = int(ih_emu * scale)
        left = int((SW - pw) / 2)
        slide.shapes.add_picture(path, left, Inches(1.9), pw, ph)
    add_slide_number(slide, num)

# ═══════════════════════════════════════════
# SLIDE 1 — COVER (full-bleed image)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK2)
cover = os.path.join(IMG, "aze.jpg")
if os.path.exists(cover):
    img = Image.open(cover)
    iw, ih = img.size
    # Fit to slide height, center horizontally (portrait image on landscape slide)
    scale = SH / (ih * 914400 / 96)
    ph = SH
    pw = int(iw * 914400 / 96 * scale)
    left = int((SW - pw) / 2)
    slide.shapes.add_picture(cover, left, 0, pw, ph)

# ═══════════════════════════════════════════
# SLIDE 2 — MÜNDƏRICAT
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)
add_title(slide, "Mündəricat", size=32)

toc_items = [
    ("1.", "Sistem icmalı", "3"),
    ("2.", "VPS quraşdırması", "4"),
    ("3.", "Obyektə quraşdırma", "5"),
    ("4.", "Klient interfeysi", "6-16"),
    ("5.", "Çoxlu kabinetlər", "17"),
    ("6.", "Yeniləmə", "18"),
    ("7.", "Diaqnostika", "19"),
]
y_start = Inches(1.2)
for i, (num, title, pg) in enumerate(toc_items):
    y = y_start + Inches(i * 0.65)
    # Number
    txBox = slide.shapes.add_textbox(Inches(2), y, Inches(0.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # Title
    txBox = slide.shapes.add_textbox(Inches(2.6), y, Inches(6), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    # Page
    txBox = slide.shapes.add_textbox(Inches(9), y, Inches(1.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = pg
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT

add_slide_number(slide, 2)

# ═══════════════════════════════════════════
# SLIDE 3 — SİSTEM İCMALI
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)
add_title(slide, "1. Sistem icmalı")
add_text(slide, (
    "Prometey — Windows maşınlarını veb-brauzer vasitəsilə\n"
    "uzaqdan idarə etmək üçün nəzərdə tutulmuş sistemdir.\n"
    "\n"
    "Üç komponent:\n"
    "• Obyekt — Windows DLL (pnpext.dll), svchost.exe yükləyir\n"
    "• VPS Relay — Python WebSocket serveri, nginx + TLS\n"
    "• Veb-klient — HTML5, istənilən brauzerdə işləyir\n"
    "\n"
    "İmkanlar: H.264/WebRTC video, fayl meneceri, terminal,\n"
    "proses/xidmət idarəsi, səs yazısı, ekran şəkilləri,\n"
    "registr, hadisə jurnalı, təhlükə monitorinqi"
), x=Inches(0.5), y=Inches(0.9), w=Inches(5.5), h=Inches(5.5), size=13)

path = os.path.join(IMG, "syn_architecture.png")
if os.path.exists(path):
    slide.shapes.add_picture(path, Inches(6.5), Inches(1.2), Inches(6.2), Inches(4.5))
add_slide_number(slide, 3)

# ═══════════════════════════════════════════
# SLIDE 4 — VPS QURAŞDIRMASI
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)
add_title(slide, "2. VPS quraşdırması")
add_text(slide, (
    "Tələblər: Ubuntu 20.04+ / Debian 11+, root, 1+ GB RAM\n"
    "\n"
    "1. Faylları VPS-ə yükləyin (scp ilə)\n"
    "2. sudo bash deploy-vps.sh işə salın\n"
    "3. Skript avtomatik: nginx, relay, coturn, TLS, portlar\n"
    "\n"
    "Nəticə:\n"
    "• Veb-klient: https://VPS_IP/\n"
    "• Obyekt: wss://VPS_IP:443/host"
), x=Inches(0.5), y=Inches(0.9), w=Inches(5), h=Inches(4), size=13)

path = os.path.join(IMG, "syn_vps_deploy.png")
if os.path.exists(path):
    slide.shapes.add_picture(path, Inches(5.8), Inches(1.0), Inches(7), Inches(5.5))
add_slide_number(slide, 4)

# ═══════════════════════════════════════════
# SLIDE 5 — OBYEKTƏ QURAŞDIRMA
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)
add_title(slide, "3. Obyektə quraşdırma (Windows)")
add_text(slide, (
    "svchost.exe ServiceDll mexanizmi (Windows standart)\n"
    "\n"
    "Fayllar: pnpext.dll, pnpext.sys (şifrələnmiş konfiq)\n"
    "Quraşdırma: install.bat Administrator kimi\n"
    "\n"
    "Yoxlama:\n"
    "• sc.exe qc WPnpSvc → svchost.exe -k PnpExtGroup\n"
    "• netstat → VPS:443-ə 2 ESTABLISHED\n"
    "\n"
    "Avto-başlanğıc: hər Windows yükləməsində"
), x=Inches(0.5), y=Inches(0.9), w=Inches(5.5), h=Inches(4), size=13)

path1 = os.path.join(IMG, "syn_host_install.png")
path2 = os.path.join(IMG, "syn_host_netstat.png")
if os.path.exists(path1):
    slide.shapes.add_picture(path1, Inches(6.5), Inches(1.0), Inches(6), Inches(2.8))
if os.path.exists(path2):
    slide.shapes.add_picture(path2, Inches(6.5), Inches(4.2), Inches(6), Inches(1.5))
add_slide_number(slide, 5)

# ═══════════════════════════════════════════
# SLIDES 6-16 — KLİENT İNTERFEYSİ
# ═══════════════════════════════════════════
client_slides = [
    (6, "4.1 Giriş", "01_login.png", "https://VPS_IP/ açın. Token və şifrəni daxil edin, Qoşul basın."),
    (7, "4.2 Panel", "04_dashboard.png", "RAM, CPU, GPU, İş vaxtı, Proseslər, Xidmətlər. Fəaliyyət jurnalı və tez əməliyyatlar."),
    (8, "4.3 Ekran", "06_screen.png", "Canlı masaüstü axını. H.264 WebRTC (UDP) və ya JPEG WebSocket (TCP, fallback)."),
    (9, "4.4 Fayllar", "07_files.png", "Fayl meneceri: disklər, endirmə/yükləmə, qovluq yaratma, drag-and-drop. Sürət: ~2 MB/s."),
    (10, "4.5 Proseslər", "08_processes.png", "CPU%, yaddaş, thread-lər. Sıralama, Kill, proqram başlatma (Adi/Admin/Sistem)."),
    (11, "4.6 Terminal", "09_terminal.png", "Uzaq CMD terminalı. ipconfig, systeminfo, dir, tasklist — real vaxtda."),
    (12, "4.7 Səs yazısı", "10_audio.png", "Yaz / Canlı / Hər ikisi. DSP: səs təmizləyici, normallaşdırma, uğultu filtri. Ekvalayzeri."),
    (13, "4.8 Ekran şəkilləri", "11_screenshots.png", "Avto-çəkmə: interval, keyfiyyət, miqyas. Tətbiqə görə filtr. Şəbəkə/Siyahı."),
    (14, "4.9 Xidmətlər", "13_services.png", "Windows xidmətləri: Başlat, Dayan, Yenidən başlat. Proqramlar: CSV ixrac."),
    (15, "4.10 Tənzimləmələr", "16_settings_update.png", "Uzaq yeniləmə, VPS Deploy, Konfiq redaktoru, Təhlükə monitoru, Özü-məhv."),
    (16, "4.11 Dil seçicisi", "20_lang_switch.png", "EN, RU, AZ arasında ani keçid. Brauzerdə yadda saxlanılır."),
]
for num, title, fname, desc in client_slides:
    img_slide(num, title, fname, desc)

# ═══════════════════════════════════════════
# SLIDE 17 — ÇOXLU KABİNETLƏR
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)
add_title(slide, "5. Çoxlu kabinetlər")
add_text(slide, (
    "Bir VPS 100-ə qədər təcrid olunmuş otağa xidmət göstərir.\n"
    "\n"
    "Hər kabinet = unikal token + şifrə\n"
    "Otaqlar avtomatik yaradılır\n"
    "Məlumatlar təcrid olunur (şəkillər, səs, fayllar)\n"
    "\n"
    "Vacib: Eyni tokeni fərqli obyektlər üçün istifadə etməyin!"
), x=Inches(1), y=Inches(1.5), w=Inches(10), h=Inches(4), size=18)

path = os.path.join(IMG, "syn_config.png")
if os.path.exists(path):
    slide.shapes.add_picture(path, Inches(8), Inches(3.5), Inches(4.5), Inches(3))
add_slide_number(slide, 17)

# ═══════════════════════════════════════════
# SLIDE 18 — YENİLƏMƏ
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)
add_title(slide, "6. Yeniləmə")
add_text(slide, (
    "Brauzerdən:\n"
    "Tənzimləmələr → VPS Deploy → Hamısını yüklə → Yenidən başlat\n"
    "\n"
    "DLL yeniləmə:\n"
    "Tənzimləmələr → Yeniləmə → URL → Yenilə\n"
    "7 addım: Endir → Defender → Dayan → Əvəz → Başlat → Defender → Tamamlandı\n"
    "\n"
    "Konfiq (yenidən quraşdırmadan):\n"
    "Tənzimləmələr → Konfiq redaktoru → Yüklə → Redaktə → Saxla (hot-reload)\n"
    "\n"
    "Yenidən başlat: Tənzimləmələr → ↻ (~5 saniyə)"
), x=Inches(1), y=Inches(1.2), w=Inches(10), h=Inches(5), size=16)
add_slide_number(slide, 18)

# ═══════════════════════════════════════════
# SLIDE 19 — DİAQNOSTİKA
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_logo(slide)
add_title(slide, "7. Diaqnostika")

problems = [
    ("Obyekt qoşulmur", "port: 443, use_tls: true yoxlayın. netstat → ESTABLISHED"),
    ("Auth failed", "Token/şifrə uyğunsuzluğu. 3 uğursuzluqdan sonra 5 dəq gözləmə"),
    ("Axın başlamır", "STUN/TURN tənzimləmələrini yoxlayın. WebRTC söndürün → JPEG fallback"),
    ("Yavaş fayllar", "Normal 1.5-2.5 MB/s TLS ilə. Çox klient kanala bölüşür"),
    ("Yüksək CPU", "Threat Monitor / Event Log cleaner söndürün. Streaming zamanı normaldır"),
]

y = Inches(1.2)
for prob, solution in problems:
    # Problem label (cyan box)
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(3), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "⚠ " + prob
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    # Solution
    txBox = slide.shapes.add_textbox(Inches(4), y, Inches(8.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "→ " + solution
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    y += Inches(1.1)

add_slide_number(slide, 19)

# ═══════════════════════════════════════════
# SLIDE 20 — FINAL
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK2)

# Logo centered, larger
logo = os.path.join(IMG, "logo.jpg")
if os.path.exists(logo):
    img = Image.open(logo)
    aspect = img.size[0] / img.size[1]
    h = Inches(2.5)
    w = int(h * aspect)
    left = (SW - w) // 2
    slide.shapes.add_picture(logo, left, Inches(1.2), w, h)

# Title
txBox = slide.shapes.add_textbox(Inches(0), Inches(4.2), SW, Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Prometey v1.0.95"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox = slide.shapes.add_textbox(Inches(0), Inches(5.2), SW, Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Uzaq İdarəetmə Sistemi"
p.font.size = Pt(18)
p.font.color.rgb = MUTED
p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
prs.save(OUT)
print(f"Created: {OUT} ({len(prs.slides)} slides)")
