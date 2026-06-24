import sys
sys.path.insert(0, r'C:\Users\Test\AppData\Roaming\Python\Python311\site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── Paths ──────────────────────────────────────────────────────────────────
DOCS  = r'D:\Android_Projects\NEW_RDP_Cloud\docs'
IMG   = os.path.join(DOCS, 'img')
OUT   = os.path.join(DOCS, 'Prometey_Presentation.pptx')

def img(name): return os.path.join(IMG, name)

# ─── Colors ─────────────────────────────────────────────────────────────────
BG_CONTENT = RGBColor(0x1A, 0x1A, 0x2E)   # very dark navy   – content slides
BG_TITLE   = RGBColor(0x2E, 0x40, 0x57)   # dark blue-gray   – title / end slides
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED    = RGBColor(0x88, 0x88, 0x88)
C_LANG     = RGBColor(0x4F, 0xC3, 0xF7)   # light blue / cyan for AZ/RU/EN labels

# ─── Presentation setup ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
blank_layout = prs.slide_layouts[6]  # blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

# ─── Background helper ───────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ─── Text box helper ─────────────────────────────────────────────────────────
def tb(slide, text, x, y, w, h, size=12, bold=False,
       color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.color.rgb  = color
    run.font.name       = 'Arial'
    return box

# ─── Multi-paragraph text box ────────────────────────────────────────────────
# lines: list of (text, size, bold, color, align)
def tb_ml(slide, lines, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(size)
        r.font.bold      = bold
        r.font.color.rgb = color
        r.font.name      = 'Arial'
    return box

# ─── Add image helper ────────────────────────────────────────────────────────
def add_img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# ─── Small logo top-left (content slides 2-19) ───────────────────────────────
def add_logo(slide):
    add_img(slide, img('logo.jpg'), 0.30, 0.20, 0.34, 0.50)

# ─── Slide number bottom-right ───────────────────────────────────────────────
def add_slide_num(slide, num):
    tb(slide, str(num), 12.60, 7.05, 0.68, 0.35,
       size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)

# ─── Standard left-text / right-image layout ─────────────────────────────────
# Left text box:  x=0.30, y=0.70, w=5.60, h=6.00
# Right image:    x=6.00, y=0.90, w=6.90, h=5.50
LT_X, LT_Y, LT_W, LT_H = 0.30, 0.70, 5.60, 6.00
RI_X, RI_Y, RI_W, RI_H = 6.00, 0.90, 6.90, 5.50

# ─── Title style for content slides ─────────────────────────────────────────
def add_title(slide, text):
    tb(slide, text, LT_X, 0.50, 5.80, 0.60, size=20, bold=True, color=C_WHITE)

# ─── Body text block ─────────────────────────────────────────────────────────
# lines: list of (text, size, bold, color)
def add_body(slide, lines, x=LT_X, y=LT_Y, w=LT_W, h=LT_H):
    tb_ml(slide, [(t, s, b, c, PP_ALIGN.LEFT) for t, s, b, c in lines],
          x, y, w, h)

# ─── Helper: language-labeled triplet ────────────────────────────────────────
def lang3(az, ru, en, size=11):
    """Return lines list for a tri-lingual block."""
    return [
        ('AZ: ' + az, size, False, C_LANG),
        ('RU: ' + ru, size, False, C_WHITE),
        ('EN: ' + en, size, False, C_MUTED),
        ('', size, False, C_WHITE),  # blank line spacer
    ]


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_TITLE)

# Large logo on right half
add_img(sl, img('logo.jpg'), 4.50, 1.00, 4.80, 5.00)

# Left half text
tb(sl, 'Data', 0.50, 1.50, 4.00, 1.20, size=52, bold=True, color=C_WHITE)
tb(sl, 'v1.0.250', 0.50, 2.80, 4.00, 0.50, size=16, color=C_MUTED)
tb_ml(sl, [
    ('Система удалённого управления и мониторинга', 14, False, C_WHITE, PP_ALIGN.LEFT),
    ('Remote Management & Monitoring System',       13, False, C_LANG,  PP_ALIGN.LEFT),
    ('Uzaqdan İdarəetmə və Monitorinq Sistemi',     13, False, C_MUTED, PP_ALIGN.LEFT),
], 0.50, 3.40, 3.90, 2.00)
add_slide_num(sl, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, 'Mündəricat / Содержание / Contents')
add_body(sl, [
    ('',                                                                            11, False, C_WHITE),
    ('1.  Sistem icmalı / Обзор системы / System Overview ............... 3',      12, False, C_WHITE),
    ('2.  VPS quraşdırması / Установка VPS / VPS Setup ................... 4',      12, False, C_WHITE),
    ('3.  Obyektə quraşdırma / Установка на хост / Host Install ......... 5',      12, False, C_WHITE),
    ('4.  Klient interfeysi / Интерфейс / Client Interface ................. 6–16', 12, False, C_WHITE),
    ('5.  Çoxlu kabinetlər / Мультикабинеты / Multi-Cabinet ............. 17',     12, False, C_WHITE),
    ('6.  Yeniləmə / Обновление / Remote Update ........................... 18',    12, False, C_WHITE),
    ('7.  Diaqnostika / Диагностика / Diagnostics ........................... 19',  12, False, C_WHITE),
], x=0.50, y=1.40, w=12.00, h=5.50)
add_slide_num(sl, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '1. Sistem icmalı / Обзор / System Overview')

lines = [
    ('AZ: Data — Windows maşınlarını veb-brauzer vasitəsilə uzaqdan idarə etmək üçün sistem',
     11, False, C_LANG),
    ('RU: Data — система удалённого управления Windows через браузер',
     11, False, C_WHITE),
    ('EN: Data — remote control system for Windows via browser',
     11, False, C_MUTED),
    ('', 6, False, C_WHITE),
    ('Components:', 11, True, C_WHITE),
    ('• Host (pnpext.dll)   • VPS Relay (Python+nginx)   • Web Client (HTML5)',
     11, False, C_WHITE),
    ('', 6, False, C_WHITE),
    ('Features:', 11, True, C_WHITE),
    ('Screen H.264/MJPEG · Files · Terminal · Processes · Audio',
     11, False, C_WHITE),
    ('Screenshots · Registry · EventLog · Services · Programs',
     11, False, C_WHITE),
]
add_body(sl, lines)
add_img(sl, img('syn_architecture.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — VPS SETUP
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '2. VPS quraşdırması / Установка VPS / VPS Setup')

lines = [
    ('Requires: Ubuntu 20.04+ / Debian, root, 1+ GB RAM', 11, True, C_WHITE),
    ('', 5, False, C_WHITE),
    ('Steps:', 11, True, C_WHITE),
    ('1. Upload files (scp)', 11, False, C_WHITE),
    ('2. Run:  sudo bash deploy-vps.sh', 11, False, C_WHITE),
    ('3. Script auto-configures nginx, relay, TLS', 11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('Result:  Web panel → https://VPS_IP/', 12, True, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Скрипт автоматически настраивает nginx, relay, TLS',
     11, False, C_LANG),
    ('AZ: Skript avtomatik nginx, relay, TLS quraşdırır',
     11, False, C_MUTED),
]
add_body(sl, lines)
add_img(sl, img('syn_vps_deploy.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOST INSTALLATION
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '3. Obyektə quraşdırma / Установка / Host Install')

lines = [
    ('Files:', 11, True, C_WHITE),
    ('pnpext.dll  ·  pnpext.sys (encrypted config)  ·  install.bat',
     11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('Run install.bat as Administrator', 12, True, C_WHITE),
    ('', 5, False, C_WHITE),
    ('Service: MspIscSvc — auto-starts on Windows boot', 11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('AZ: install.bat Administrator kimi işə salın', 11, False, C_LANG),
    ('RU: Запустить install.bat от администратора',   11, False, C_MUTED),
]
add_body(sl, lines)
add_img(sl, img('syn_host_install.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LOGIN  (4.1)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.1 Giriş / Вход / Login')

lines = [
    ('Open https://VPS_IP/ → Enter login + password → Connect',
     12, True, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Откройте https://VPS_IP/ → Введите логин и пароль → Подключиться',
     11, False, C_WHITE),
    ('AZ: https://VPS_IP/ açın → Giriş + şifrəni daxil edin → Qoşul basın',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('01_login.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DASHBOARD  (4.2)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.2 Panel / Дашборд / Dashboard')

lines = [
    ('CPU, RAM, GPU, Disk — real-time',                   12, True,  C_WHITE),
    ('VPS1 + VPS2 cards   ·   Speed test',                11, False, C_WHITE),
    ('Activity log  @username',                           11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: CPU, RAM, GPU, диск в реальном времени. Карточки VPS1+VPS2. Тест скорости.',
     11, False, C_WHITE),
    ('AZ: CPU, RAM, GPU, disk — real vaxt. VPS1+VPS2 kartları. Sürət testi.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('04_dashboard.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SCREEN STREAM  (4.3)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.3 Ekran / Экран / Screen')

lines = [
    ('Live stream H.264 / MJPEG,  60 FPS',               12, True,  C_WHITE),
    ('Action buttons: Fit/Fill · Screenshot · Audio · Record · Fullscreen',
     11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Прямая трансляция H.264/MJPEG, 60 FPS.',        11, False, C_WHITE),
    ('    Кнопки: Вписать · Скриншот · Звук · Запись · Полный экран',
     11, False, C_WHITE),
    ('AZ: Canlı axın H.264/MJPEG, 60 FPS.',               11, False, C_LANG),
    ('    Düymələr: Fit · Ekran şəkli · Səs · Yaz · Tam ekran',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('06_screen.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FILE MANAGER  (4.4)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.4 Fayllar / Файлы / Files')

lines = [
    ('Full filesystem access',            12, True,  C_WHITE),
    ('Download / Upload   ·   Drag & drop',11, False, C_WHITE),
    ('~2 MB/s',                            11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Полный доступ к файловой системе. Скачать/Загрузить. ~2 МБ/с.',
     11, False, C_WHITE),
    ('AZ: Fayl sisteminə tam giriş. Endir/Yüklə. Sürükle-burax. ~2 MB/s.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('07_files.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PROCESSES  (4.5)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.5 Proseslər / Процессы / Processes')

lines = [
    ('PID · Name · CPU% · RAM',            12, True,  C_WHITE),
    ('Kill · Start as SYSTEM',             11, False, C_WHITE),
    ('Services: Start / Stop / Restart',   11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: PID, имя, CPU%, RAM. Завершить, Запустить как SYSTEM. Службы: Старт/Стоп.',
     11, False, C_WHITE),
    ('AZ: PID, ad, CPU%, RAM. Bitir, SYSTEM kimi başlat. Xidmətlər: Başlat/Dayan.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('08_processes.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TERMINAL  (4.6)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.6 Terminal / Терминал / Terminal')

lines = [
    ('cmd + PowerShell,  runs as SYSTEM',  12, True,  C_WHITE),
    ('Real-time output',                   11, False, C_WHITE),
    ('dir  ·  tasklist  ·  Get-Process…',  11, False, C_MUTED),
    ('', 5, False, C_WHITE),
    ('RU: cmd + PowerShell от SYSTEM. Вывод в реальном времени.',
     11, False, C_WHITE),
    ('AZ: cmd + PowerShell, SYSTEM kimi. Real vaxt çıxışı.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('09_terminal.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — AUDIO  (4.7)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.7 Səs / Аудио / Audio')

lines = [
    ('Continuous recording  OGG Opus,  5-min segments',   12, True,  C_WHITE),
    ('WASAPI SYSTEM  (no privacy indicator)',               11, False, C_WHITE),
    ('DSP: denoise · hum filter',                          11, False, C_WHITE),
    ('Live listen: WASAPI loopback',                       11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Непрерывная запись OGG Opus, 5-мин сегменты. WASAPI SYSTEM (без индикатора). Живое прослушивание.',
     11, False, C_WHITE),
    ('AZ: Davamlı yazma OGG Opus, 5 dəq seqmentlər. WASAPI SYSTEM. Canlı dinləmə.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('10_audio.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SCREENSHOTS  (4.8)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.8 Ekran şəkilləri / Скриншоты / Screenshots')

lines = [
    ('Auto every 10 s  (configurable)',                12, True,  C_WHITE),
    ('Gallery with inline preview',                    11, False, C_WHITE),
    ('Filter by app  ·  Quality / scale settings',     11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Авто каждые 10 сек. Галерея с предпросмотром. Фильтр по приложению.',
     11, False, C_WHITE),
    ('AZ: Avto hər 10 san. Önizləmə ilə qalereya. Tətbiqə görə filtr.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('11_screenshots.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — EVENT LOG + REGISTRY  (4.9)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.9 Hadisə + Reyestr / Журнал + Реестр / EventLog + Registry')

lines = [
    ('Event Log:',                                          11, True,  C_WHITE),
    ('  Application · System · Security sources',          11, False, C_WHITE),
    ('Registry:',                                          11, True,  C_WHITE),
    ('  HKLM · HKCU · HKCR · HKU · HKCC — browse, edit, delete',
     11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Журнал: Application/System/Security. Реестр: Полный браузер, редактирование.',
     11, False, C_WHITE),
    ('AZ: Jurnal: Application/System/Security. Reyestr: Tam brauzer, redaktə.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('12_eventlog.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DEFENSE / THREAT  (4.10)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.10 Müdafiə / Защита / Defense')

lines = [
    ('Windows Defender management',                        12, True,  C_WHITE),
    ('Enable / Disable  ·  Scan threats',                  11, False, C_WHITE),
    ('Event log cleanup',                                  11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Управление Windows Defender. Вкл/выкл. Сканирование угроз. Очистка журнала.',
     11, False, C_WHITE),
    ('AZ: Windows Defender idarəsi. Aç/bağla. Təhdid skanı. Jurnal təmizləmə.',
     11, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('19_settings_threat.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — SETTINGS (4 tabs, 2×2 grid)  (4.11)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '4.11 Tənzimləmələr / Настройки / Settings')

lines = [
    ('4 tabs:',                                             12, True,  C_WHITE),
    ('Update — remote DLL update',                          11, False, C_WHITE),
    ('Deploy — VPS config',                                 11, False, C_WHITE),
    ('Config — agent params',                               11, False, C_WHITE),
    ('Threat — security monitor',                           11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: 4 вкладки: Обновление · Развёртывание · Конфиг · Угрозы',
     11, False, C_WHITE),
    ('AZ: 4 tab: Yeniləmə · Deploy · Konfiq · Təhdid',
     11, False, C_LANG),
]
add_body(sl, lines)

# 2×2 grid of settings screenshots on right
grid_imgs = [
    ('16_settings_update.png', 6.10, 0.90),
    ('17_settings_deploy.png', 9.75, 0.90),
    ('18_settings_config.png', 6.10, 3.80),
    ('19_settings_threat.png', 9.75, 3.80),
]
for fname, gx, gy in grid_imgs:
    add_img(sl, img(fname), gx, gy, 3.45, 2.70)
add_slide_num(sl, 16)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — MULTI-CABINET / USERS
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '5. Çoxlu kabinetlər / Мультикабинеты / Multi-Cabinet')

lines = [
    ('One VPS serves 100+ isolated rooms',                  12, True,  C_WHITE),
    ('Admin role: all features',                            11, False, C_WHITE),
    ('Operator role: permitted tabs only',                  11, False, C_WHITE),
    ('Each operator @username logged in Activity Log',      11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: Один VPS — 100+ кабинетов. Admin: полный доступ. Operator: только разрешённые вкладки. @имя в журнале.',
     11, False, C_WHITE),
    ('AZ: Bir VPS — 100+ kabinet. Admin: tam giriş. Operator: icazəli tablar. @ad jurnalda.',
     11, False, C_LANG),
]
add_body(sl, lines, x=0.30, y=1.40, w=12.50, h=5.50)
add_slide_num(sl, 17)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — REMOTE UPDATE
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '6. Yeniləmə / Обновление / Remote Update')

lines = [
    ('5 steps:',                                            11, True,  C_WHITE),
    ('1. Upload pnpext.dll to VPS',                         11, False, C_WHITE),
    ('2. Click "Update Agent"',                             11, False, C_WHITE),
    ('3. Agent downloads  ~15–30 s',                        11, False, C_WHITE),
    ('4. Brief disconnect  ~5 s',                           11, False, C_WHITE),
    ('5. Status Online — done.  No physical access needed.', 11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('RU: 1. Загрузить pnpext.dll  2. Нажать "Обновить агент"  3. ~15-30с  4. ~5с разрыв  5. Online.',
     10, False, C_WHITE),
    ('AZ: 1. pnpext.dll yüklə  2. "Agenti yenilə" bas  3. ~15-30s  4. ~5s kəsilmə  5. Online.',
     10, False, C_LANG),
]
add_body(sl, lines)
add_img(sl, img('16_settings_update.png'), RI_X, RI_Y, RI_W, RI_H)
add_slide_num(sl, 18)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — DIAGNOSTICS (text only, no screenshot)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_CONTENT)
add_logo(sl)
add_title(sl, '7. Diaqnostika / Диагностика / Diagnostics')

lines = [
    ('AZ:', 11, True, C_LANG),
    ('⚠ Qoşulmur → port 443, use_tls: true yoxlayın',               11, False, C_LANG),
    ('⚠ Auth failed → Token/şifrə yanlışdır. 3 cəhd → 5 dəq gözləmə', 11, False, C_LANG),
    ('⚠ Axın başlamır → STUN/TURN yoxlayın, WebRTC söndürün → MJPEG', 11, False, C_LANG),
    ('⚠ Yavaş fayllar → Normal 1.5–2.5 MB/s TLS ilə',               11, False, C_LANG),
    ('', 5, False, C_WHITE),
    ('RU:', 11, True, C_WHITE),
    ('⚠ Не подключается → проверьте port 443, use_tls: true',         11, False, C_WHITE),
    ('⚠ Auth failed → Неверный логин/пароль. 3 попытки → пауза 5 мин', 11, False, C_WHITE),
    ('⚠ Стрим не стартует → проверьте STUN/TURN, отключите WebRTC → MJPEG', 11, False, C_WHITE),
    ('⚠ Медленные файлы → Норма 1.5–2.5 MB/s с TLS',                 11, False, C_WHITE),
    ('', 5, False, C_WHITE),
    ('EN:', 11, True, C_MUTED),
    ('⚠ Not connecting → check port 443, use_tls: true',              11, False, C_MUTED),
    ('⚠ Auth failed → Wrong credentials. 3 failures → 5 min pause',  11, False, C_MUTED),
    ('⚠ Stream won\'t start → check STUN/TURN, disable WebRTC → MJPEG fallback', 11, False, C_MUTED),
    ('⚠ Slow files → Normal 1.5–2.5 MB/s with TLS',                  11, False, C_MUTED),
]
add_body(sl, lines, x=0.30, y=1.30, w=12.50, h=6.00)
add_slide_num(sl, 19)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — END / THANK YOU
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
set_bg(sl, BG_TITLE)

# Large centered logo
add_img(sl, img('logo.jpg'), 5.665, 1.30, 2.00, 2.00)

tb_ml(sl, [
    ('Data v1.0.250',  24, True,  C_WHITE, PP_ALIGN.CENTER),
], 0.50, 3.60, 12.33, 0.70)

tb_ml(sl, [
    ('Uzaq İdarəetmə Sistemi',                        13, False, C_LANG,  PP_ALIGN.CENTER),
    ('Система удалённого управления',                 13, False, C_WHITE, PP_ALIGN.CENTER),
    ('Remote Management System',                      13, False, C_MUTED, PP_ALIGN.CENTER),
], 0.50, 4.45, 12.33, 1.60)
add_slide_num(sl, 20)


# ─── Save ────────────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
size = os.path.getsize(OUT)
print(f'Saved:  {OUT}')
print(f'Slides: {len(prs.slides)}')
print(f'Size:   {size:,} bytes')
