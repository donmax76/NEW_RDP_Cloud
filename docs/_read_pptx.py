import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Pt

fname = sys.argv[1] if len(sys.argv) > 1 else r'D:\Android_Projects\NEW_RDP_Cloud\docs\Prometey_AZ.pptx'
prs = Presentation(fname)
print(f'Slides: {len(prs.slides)}')
print(f'Size: {prs.slide_width.inches:.2f}" x {prs.slide_height.inches:.2f}"')
print()
for i, slide in enumerate(prs.slides, 1):
    print(f'--- SLIDE {i} ---')
    bg = slide.background.fill
    try:
        if bg.type is not None:
            try:
                print(f'  BG: #{bg.fore_color.rgb}')
            except:
                pass
    except:
        pass
    for shape in slide.shapes:
        if not shape.has_text_frame:
            if shape.shape_type == 13:
                print(f'  [IMAGE] {shape.name} pos=({shape.left/914400:.2f}", {shape.top/914400:.2f}") size=({shape.width/914400:.2f}"x{shape.height/914400:.2f}")')
            continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            info = []
            for run in para.runs:
                if run.font.size:
                    info.append(f'{int(run.font.size.pt)}pt')
                if run.font.bold:
                    info.append('BOLD')
                try:
                    if run.font.color.type is not None:
                        info.append(f'#{run.font.color.rgb}')
                except:
                    pass
                break
            tag = '|'.join(dict.fromkeys(info)) if info else '?'
            print(f'  [{tag}] {t[:120]}')
    print()
